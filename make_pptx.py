#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Soutenance PFE - Barrage Ahl Souss - Simulation OpenFOAM"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as RGB
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

BASE = "/workspace"
FIG  = BASE + "/figures/"
LOG  = BASE + "/logos/"
OUT  = BASE + "/Presentation_Soutenance_BarrageAhlSouss.pptx"

# ── Colors ────────────────────────────────────────────
DARK  = RGB(0x02,0x15,0x26)
DARK2 = RGB(0x0A,0x25,0x40)
PRIM  = RGB(0x06,0x5A,0x82)
TEAL  = RGB(0x1C,0x72,0x93)
TEAL2 = RGB(0xD4,0xEE,0xF7)
ACCN  = RGB(0xE8,0xA8,0x38)
LITE  = RGB(0xEB,0xF5,0xFB)
WITE  = RGB(0xFF,0xFF,0xFF)
DT    = RGB(0x1A,0x2A,0x3A)
MT    = RGB(0x4A,0x6B,0x7A)
LT    = RGB(0xC8,0xDF,0xE8)
GREEN = RGB(0x1A,0x7A,0x3A)
ORANG = RGB(0xC4,0x7A,0x00)
RED   = RGB(0xC0,0x39,0x2B)

prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(5.625)
BLANK = prs.slide_layouts[6]

# ── Helpers ───────────────────────────────────────────

def ns():
    return prs.slides.add_slide(BLANK)

def bg(s, c):
    f = s.background.fill; f.solid(); f.fore_color.rgb = c

def rect(s, x, y, w, h, fill=None, line=None, rnd=False, lw=0.75):
    sh = s.shapes.add_shape(5 if rnd else 1,
                            Inches(x), Inches(y), Inches(w), Inches(h))
    if fill: sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:    sh.fill.background()
    if line: sh.line.color.rgb = line; sh.line.width = Pt(lw)
    else:    sh.line.fill.background()
    return sh

def ov(s, x, y, w, h, fill):
    sh = s.shapes.add_shape(9, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    return sh

def tx(s, text, x, y, w, h, sz=11, fc=None, bold=False,
       italic=False, al=PP_ALIGN.LEFT, fn="Calibri"):
    if fc is None: fc = DT
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(str(text).split('\n')):
        pg = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        pg.alignment = al
        rn = pg.add_run(); rn.text = line
        rn.font.size = Pt(sz); rn.font.bold = bold
        rn.font.italic = italic; rn.font.name = fn
        rn.font.color.rgb = fc

def im(s, path, x, y, w, h):
    try: s.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    except Exception as e: print(f"  IMG: {path}: {e}")

def note(s, text):
    s.notes_slide.notes_text_frame.text = text

def card(s, x, y, w, h):
    return rect(s, x, y, w, h, fill=WITE, line=TEAL2, rnd=True)

def tcell(cell, text, bg_c=None, fc=DT, sz=10, bold=False,
          al=PP_ALIGN.CENTER):
    if bg_c: cell.fill.solid(); cell.fill.fore_color.rgb = bg_c
    tf = cell.text_frame; pg = tf.paragraphs[0]; pg.alignment = al
    if pg.runs: rn = pg.runs[0]
    else:       rn = pg.add_run()
    rn.text = text; rn.font.size = Pt(sz); rn.font.bold = bold
    rn.font.color.rgb = fc; rn.font.name = "Calibri"

def lslide(title, sub=None):
    s = ns(); bg(s, LITE)
    tx(s, title, 0.4, 0.1, 9.2, 0.62, sz=23, fc=PRIM, bold=True, fn="Cambria")
    if sub: tx(s, sub, 0.4, 0.68, 9.2, 0.28, sz=9.5, fc=MT, italic=True)
    return s

def dslide(title, sub=None):
    s = ns(); bg(s, DARK)
    tx(s, title, 0.4, 0.1, 9.2, 0.58, sz=21, fc=WITE, bold=True, fn="Cambria")
    if sub: tx(s, sub, 0.4, 0.64, 9.2, 0.28, sz=9.5, fc=LT, italic=True)
    return s

# ══ SLIDE 1 — TITRE ══════════════════════════════════
s = ns(); bg(s, DARK)
im(s, LOG+"ehtp-logo.png", 0.3, 0.14, 1.6, 0.85)
im(s, LOG+"cid-logo.jpg",  8.1, 0.14, 1.6, 0.85)
tx(s, "École Hassania des Travaux Publics  —  Génie de l'Eau et de l'Environnement",
   0.5, 0.2, 9, 0.3, sz=9, fc=LT, al=PP_ALIGN.CENTER)
im(s, FIG+"figure4_location.png", 6.0, 0.9, 3.7, 4.5)
tx(s, "SIMULATION HYDRAULIQUE",   0.45, 1.22, 7, 0.75, sz=25, fc=WITE, bold=True, fn="Cambria")
tx(s, "AVEC OpenFOAM",            0.45, 1.95, 7, 0.62, sz=23, fc=ACCN, bold=True, fn="Cambria")
tx(s, "Barrage Ahl Souss (Aït M'zal)",
   0.45, 2.65, 5.5, 0.48, sz=16, fc=LT, italic=True, fn="Cambria")
tx(s, "Dimensionnement et validation numérique\nde l'évacuateur de crues  —  CFD 3D",
   0.45, 3.18, 5.5, 0.65, sz=11.5, fc=LT)
rect(s, 0.4, 3.95, 5.5, 0.58, fill=PRIM, line=TEAL, rnd=True)
tx(s, "Hiba BASTOUT   •   Asmae BOULEMKHARJ",
   0.45, 3.98, 5.4, 0.52, sz=13, fc=WITE, bold=True, al=PP_ALIGN.CENTER)
tx(s, "Encadrants :  Mme AHATTAB Jihane  |  Mme BENABDELLAH Khaoula  |  M. SOULI Hamza",
   0.45, 4.63, 9, 0.3, sz=8.5, fc=LT)
tx(s, "Juin 2026", 0.45, 4.95, 3, 0.25, sz=9, fc=LT, italic=True)
note(s, "Slide 1 - Titre. Présentez auteurs, école, encadrants. (30 sec)")

# ══ SLIDE 2 — PLAN ═══════════════════════════════════
s = lslide("Plan de la présentation")
cards_data = [
    ("01", PRIM,  "Contexte & Site",
     "Barrage Ahl Souss — localisation, caractéristiques, débits de projet"),
    ("02", TEAL,  "Dimensionnement théorique",
     "Seuil Creager, coursier, cuillère saut de ski — méthode analytique 1D"),
    ("03", ACCN,  "Modélisation OpenFOAM",
     "Simulation 3D diphasique — interFoam, VOF, k-ε, maillage, convergence"),
    ("04", DARK2, "Résultats & Validation",
     "Comparaison théorie/simulation — Q10, Q100, Q1000 — Simulation exploratoire"),
]
gpos = [(0.3,0.9),(5.2,0.9),(0.3,3.1),(5.2,3.1)]
for (n, col, title, desc), (cx, cy) in zip(cards_data, gpos):
    card(s, cx, cy, 4.5, 2.0)
    ov(s, cx+0.15, cy+0.15, 0.75, 0.75, col)
    tx(s, n, cx+0.15, cy+0.17, 0.75, 0.72, sz=18, fc=WITE, bold=True,
       al=PP_ALIGN.CENTER, fn="Cambria")
    tx(s, title, cx+1.05, cy+0.18, 3.3, 0.5, sz=14, fc=DT, bold=True, fn="Cambria")
    rect(s, cx+1.05, cy+0.75, 3.2, 0.03, fill=col)
    tx(s, desc, cx+0.2, cy+0.88, 4.1, 1.0, sz=10.5, fc=MT)
note(s, "Slide 2 - Plan. 4 parties. (30 sec)")

# ══ SLIDE 3 — CONTEXTE ════════════════════════════════
s = lslide("Contexte & Problématique")
im(s, FIG+"figure4_location.png", 0.3, 0.88, 4.2, 4.45)
tx(s, "Vue satellite — Barrage Ahl Souss\n(Source : Google Maps, 2024)",
   0.3, 5.05, 4.2, 0.38, sz=8, fc=MT, italic=True, al=PP_ALIGN.CENTER)
items = [
    ("Site",          "Oued Izig — Province de Chtouka Aït Baha, région Souss-Massa"),
    ("Ouvrage",       "Barrage BCR — Bassin versant : 187 km²  |  Alt. moy. : 1 287 m"),
    ("Rôles",         "Irrigation, alimentation en eau potable (AEP), protection crues"),
    ("Problématique", "Dimensionner et valider l'évacuateur de crues pour Q10, Q100, Q1000"),
    ("Enjeu",         "L'approche analytique 1D est insuffisante — validation CFD 3D nécessaire"),
]
cy = 0.9
for label, content in items:
    card(s, 4.72, cy, 5.0, 0.72)
    tx(s, label+" :", 4.88, cy+0.08, 1.5, 0.55, sz=10, fc=PRIM, bold=True)
    tx(s, content,   6.28, cy+0.08, 3.32, 0.55, sz=9.5, fc=DT)
    cy += 0.82
note(s, "Slide 3 - Contexte et problématique. (1min30)")

# ══ SLIDE 4 — DONNÉES DU SITE ════════════════════════
s = lslide("Barrage Ahl Souss — Données de base")
card(s, 0.3, 0.88, 4.5, 4.6)
tx(s, "Caractéristiques du barrage", 0.45, 0.95, 4.2, 0.42,
   sz=13, fc=PRIM, bold=True, fn="Cambria")
rect(s, 0.45, 1.4, 4.0, 0.03, fill=PRIM)
carac = [
    ("Type",          "Béton Compacté au Rouleau (BCR)"),
    ("Cours d'eau",   "Oued Izig"),
    ("Province",      "Chtouka Aït Baha"),
    ("Bassin vers.",  "187 km²"),
    ("Altitude moy.", "1 287 m NGM"),
    ("Évacuateur",    "Seuil Creager + coursier + cuillère"),
]
ry = 1.52
for k, v in carac:
    tx(s, k, 0.45, ry, 2.0, 0.42, sz=10.5, fc=DT, bold=True)
    tx(s, v, 2.45, ry, 2.2, 0.42, sz=10.5, fc=MT)
    ry += 0.47

card(s, 5.0, 0.88, 4.7, 4.6)
tx(s, "Débits de projet", 5.15, 0.95, 4.4, 0.42, sz=13, fc=TEAL, bold=True, fn="Cambria")
rect(s, 5.15, 1.4, 4.2, 0.03, fill=TEAL)
tbl = s.shapes.add_table(4, 3, Inches(5.1), Inches(1.52), Inches(4.6), Inches(2.2)).table
for j, h in enumerate(["Période (ans)", "Q laminé (m³/s)", "Rôle"]):
    tcell(tbl.cell(0,j), h, bg_c=PRIM, fc=WITE, sz=10.5, bold=True)
for i, (per, q, role) in enumerate([("10","160,74","Vérification"),
                                     ("100","507,54","Dimensionnement"),
                                     ("1 000","794,93","Sécurité")]):
    bg_r = LITE if i%2==0 else WITE
    for j, v in enumerate([per, q, role]):
        tcell(tbl.cell(i+1,j), v, bg_c=bg_r, fc=DT, sz=10.5)
rect(s, 5.1, 3.85, 4.6, 1.45, fill=PRIM, rnd=True)
tx(s, "Débit de référence", 5.2, 3.92, 4.4, 0.35, sz=9.5, fc=LT,
   al=PP_ALIGN.CENTER, bold=True)
tx(s, "Q₁₀₀₀ = 794,93 m³/s", 5.2, 4.28, 4.4, 0.52, sz=20, fc=ACCN,
   bold=True, al=PP_ALIGN.CENTER, fn="Cambria")
tx(s, "Débit millénal laminé — cas le plus défavorable", 5.2, 4.82, 4.4, 0.32,
   sz=8.5, fc=LT, italic=True, al=PP_ALIGN.CENTER)
note(s, "Slide 4 - Données du site. (1 min)")

# ══ SLIDE 5 — OBJECTIFS ══════════════════════════════
s = lslide("Objectifs du PFE")
objs = [
    (PRIM,"1","Dimensionner théoriquement l'évacuateur de crues",
     "Calculer les dimensions du seuil Creager, du coursier et de la cuillère saut de ski "
     "pour Q10, Q100, Q1000 par les méthodes analytiques classiques (Creager, HEC-RAS)."),
    (TEAL,"2","Modéliser numériquement l'écoulement avec OpenFOAM",
     "Construire un modèle CFD 3D diphasique eau/air avec le solveur interFoam "
     "et le modèle de turbulence k-ε. Simuler les trois débits de projet."),
    (ACCN,"3","Valider les résultats et analyser les écarts",
     "Comparer les hauteurs d'eau simulées et théoriques, calculer les écarts "
     "relatifs et interpréter les différences pour évaluer la fiabilité du modèle."),
]
for i, (col, n, title, text) in enumerate(objs):
    oy = 0.88 + i*1.55
    card(s, 0.3, oy, 9.4, 1.38)
    ov(s, 0.5, oy+0.32, 0.75, 0.75, col)
    tx(s, n, 0.5, oy+0.32, 0.75, 0.75, sz=22, fc=WITE, bold=True,
       al=PP_ALIGN.CENTER, fn="Cambria")
    tx(s, title, 1.45, oy+0.1, 8.1, 0.45, sz=13, fc=DT, bold=True, fn="Cambria")
    tx(s, text,  1.45, oy+0.58, 8.1, 0.72, sz=10.5, fc=MT)
note(s, "Slide 5 - Objectifs. (45 sec)")

# ══ SLIDE 6 — MÉTHODOLOGIE ═══════════════════════════
s = lslide("Approche méthodologique")
steps = [
    (PRIM,"1","Hydrologie",       "Q10, Q100\nQ1000"),
    (TEAL,"2","Dimensionnement", "Creager\ncoursier\ncuillère"),
    (PRIM,"3","Maillage 3D",     "snappyHex\ngrossier\nraffiné"),
    (TEAL,"4","Simulation\nOpenFOAM","interFoam\nk-ε / VOF"),
    (ACCN,"5","Résultats &\nValidation","ParaView\ncomp. théo"),
]
bw, bh = 1.52, 2.22
gy = 1.25
totalW = len(steps)*bw + (len(steps)-1)*0.38
bx0 = (10 - totalW) / 2
for i, (col, n, title, sub) in enumerate(steps):
    bx = bx0 + i*(bw+0.38)
    rect(s, bx, gy, bw, bh, fill=col, rnd=True)
    tx(s, n,     bx, gy+0.06, bw, 0.42, sz=22, fc=WITE, bold=True,
       al=PP_ALIGN.CENTER, fn="Cambria")
    tx(s, title, bx+0.06, gy+0.52, bw-0.12, 0.72, sz=11, fc=WITE, bold=True,
       al=PP_ALIGN.CENTER, fn="Cambria")
    tx(s, sub,   bx+0.06, gy+1.28, bw-0.12, 0.82, sz=9, fc=LITE,
       al=PP_ALIGN.CENTER)
    if i < len(steps)-1:
        ax = bx+bw+0.04
        rect(s, ax, gy+bh/2-0.03, 0.3, 0.06, fill=MT)
        tx(s, "▶", ax+0.12, gy+bh/2-0.2, 0.22, 0.4, sz=12, fc=DT,
           al=PP_ALIGN.CENTER)
rect(s, 0.3, 3.66, 4.5, 0.48, fill=LITE, line=PRIM, rnd=True)
tx(s, "Analytique 1D : loi de Creager, courbe de remous (HEC-RAS)",
   0.42, 3.71, 4.26, 0.42, sz=9, fc=PRIM)
rect(s, 5.2, 3.66, 4.5, 0.48, fill=LITE, line=TEAL, rnd=True)
tx(s, "CFD 3D : maillage hexaédrique, solveur interFoam, turbulence k-ε RAS",
   5.32, 3.71, 4.26, 0.42, sz=9, fc=TEAL)
note(s, "Slide 6 - Méthodologie. Double approche analytique/numérique. (45 sec)")

# ══ SLIDE 7 — DIMENSIONNEMENT ════════════════════════
s = lslide("Dimensionnement de l'évacuateur de crues",
           "Méthode analytique — Seuil Creager + Coursier + Cuillère saut de ski")
card(s, 0.3, 0.88, 4.6, 4.62)
im(s, FIG+"ch8_profil_clair.jpeg", 0.38, 0.96, 4.44, 4.28)
tx(s, "Profil longitudinal de l'évacuateur",
   0.3, 5.06, 4.6, 0.32, sz=8.5, fc=MT, italic=True, al=PP_ALIGN.CENTER)
structs = [
    (PRIM,"1","Seuil Creager","Loi : Q = Cd · L · Hd^(3/2)\nHd_théo = 2,36 m  (Q₁₀₀₀ = 794,93 m³/s)"),
    (TEAL,"2","Coursier","Calcul pas-à-pas HEC-RAS\nh_théo = 0,82 m"),
    (ACCN,"3","Cuillère saut de ski","R = 3 m  |  y₀ = 0,61 m\nVérif. : R/y₀ > 4"),
]
sy = 0.9
for col, n, title, text in structs:
    card(s, 5.1, sy, 4.6, 1.12)
    ov(s, 5.2, sy+0.25, 0.6, 0.6, col)
    tx(s, n, 5.2, sy+0.25, 0.6, 0.6, sz=15, fc=WITE, bold=True,
       al=PP_ALIGN.CENTER, fn="Cambria")
    tx(s, title, 5.92, sy+0.08, 3.68, 0.4, sz=12, fc=DT, bold=True, fn="Cambria")
    tx(s, text,  5.92, sy+0.52, 3.68, 0.52, sz=10, fc=MT)
    sy += 1.22
rect(s, 5.1, 4.68, 4.6, 0.8, fill=PRIM, rnd=True)
tx(s, "Vérification R/y₀ > 4 — validée pour les 3 débits",
   5.2, 4.73, 4.4, 0.35, sz=11, fc=WITE, bold=True,
   al=PP_ALIGN.CENTER, fn="Cambria")
tx(s, "Q10 : 23,4 ✓     Q100 : 6,96 ✓     Q1000 : 4,73 ✓",
   5.2, 5.09, 4.4, 0.3, sz=10.5, fc=ACCN, bold=True, al=PP_ALIGN.CENTER)
note(s, "Slide 7 - Dimensionnement. 3 composantes. Vérification R/y0 > 4. (1min30)")

# ══ SLIDE 8 — OPENFOAM ═══════════════════════════════
s = lslide("Configuration numérique — OpenFOAM",
           "Simulation diphasique instationnaire — solveur interFoam")
cfgs = [
    (PRIM,"Solveur",["interFoam — diphasique instationnaire",
                     "Méthode VOF — champ α (0=air, 1=eau)",
                     "Équations RANS Reynolds-moyennées"]),
    (TEAL,"Turbulence",["Modèle k-ε RAS",
                        "kqRWallFunction aux parois",
                        "Résolution jusqu'à la paroi"]),
    (ACCN,"Maillage",["snappyHexMesh — hexaédrique",
                      "Grossier : cellules 70 cm (calibration)",
                      "Raffiné : cellules 20 cm (production)"]),
    (DARK2,"Convergence",["Résidus p_rgh, U < 10⁻⁵",
                          "Nombre de Courant Co < 1",
                          "Régime quasi-permanent atteint"]),
]
gpos = [(0.3,0.88),(2.72,0.88),(0.3,2.94),(2.72,2.94)]
for (col, title, items_), (cx, cy) in zip(cfgs, gpos):
    card(s, cx, cy, 2.3, 1.88)
    rect(s, cx+0.1, cy+0.1, 2.1, 0.42, fill=col, rnd=True)
    tx(s, title, cx+0.1, cy+0.1, 2.1, 0.42, sz=12, fc=WITE, bold=True,
       al=PP_ALIGN.CENTER, fn="Cambria")
    for j, it in enumerate(items_):
        tx(s, "• "+it, cx+0.15, cy+0.58+j*0.41, 2.05, 0.38, sz=9.5, fc=DT)
card(s, 5.24, 0.88, 4.46, 2.08)
im(s, FIG+"ch8_maillage_raffine.jpeg", 5.32, 0.96, 4.3, 1.85)
tx(s, "Maillage raffiné — cellules 20 cm",
   5.24, 3.02, 4.46, 0.28, sz=8.5, fc=MT, italic=True, al=PP_ALIGN.CENTER)
card(s, 5.24, 3.35, 4.46, 2.12)
tx(s, "Conditions aux limites", 5.35, 3.42, 4.25, 0.4,
   sz=12, fc=TEAL, bold=True, fn="Cambria")
bc = s.shapes.add_table(5, 3, Inches(5.3), Inches(3.88), Inches(4.35), Inches(1.47)).table
for j, h in enumerate(["Frontière","Vitesse U","α"]):
    tcell(bc.cell(0,j), h, bg_c=TEAL, fc=WITE, sz=9.5, bold=True)
for i, (fr, u, a) in enumerate([("Entrée amont","débit imposé","α = 1"),
                                  ("Sortie aval","grad. nul","grad. nul"),
                                  ("Parois","no-slip","—"),
                                  ("Atmosphère","inlet/outlet","α = 0")]):
    bg_r = LITE if i%2==0 else WITE
    for j, v in enumerate([fr, u, a]):
        tcell(bc.cell(i+1,j), v, bg_c=bg_r, fc=DT, sz=9)
note(s, "Slide 8 - OpenFOAM. interFoam, k-ε, snappyHexMesh. (2 min)")

# ══ SLIDE 9 — RÉSULTATS Q1000 VIZ ════════════════════
s = dslide("Résultats de simulation — Q₁₀₀₀ = 794,93 m³/s",
           "Champs extraits en régime quasi-permanent — visualisation ParaView")
viz = [
    (FIG+"ch9_q1000_U_profil.jpeg",  "Champ de vitesse ‖U‖ (m/s)"),
    (FIG+"ch9_q1000_P_profil.jpeg",  "Champ de pression p (Pa)"),
    (FIG+"ch8_alpha_20_raffine.png",  "Interface eau/air — champ α"),
    (FIG+"ch9_q1000_jet.png",         "Trajectoire du jet — cuillère"),
]
iw, ih = 4.55, 2.06
vpos = [(0.25,0.96),(5.2,0.96),(0.25,3.27),(5.2,3.27)]
for (path, label), (px, py) in zip(viz, vpos):
    rect(s, px, py, iw, ih+0.28, fill=DARK2, line=TEAL, rnd=True)
    im(s, path, px+0.06, py+0.06, iw-0.12, ih)
    tx(s, label, px+0.06, py+ih+0.08, iw-0.12, 0.22,
       sz=8.5, fc=LT, italic=True, al=PP_ALIGN.CENTER)
note(s, "Slide 9 - Visualisations Q1000: vitesse, pression, α, jet. (2 min)")

# ══ SLIDE 10 — COMPARAISON ════════════════════════════
s = lslide("Comparaison Théorie vs Simulation — Q₁₀₀₀")
card(s, 0.3, 0.88, 5.7, 3.5)
tx(s, "Tableau comparatif", 0.45, 0.95, 5.4, 0.38, sz=13, fc=PRIM, bold=True, fn="Cambria")
ct = s.shapes.add_table(4, 4, Inches(0.38), Inches(1.4), Inches(5.55), Inches(2.2)).table
for j, h in enumerate(["Grandeur","Théorique","Simulée","Écart (%)"]):
    tcell(ct.cell(0,j), h, bg_c=PRIM, fc=WITE, sz=10.5, bold=True)
comp = [("Hd — Seuil Creager (m)","2,36","2,54","7,6 %",GREEN),
        ("h — Coursier (m)",      "0,82","1,12","36,6 %",ORANG),
        ("h — Cuillère (m)",      "0,61","1,01","65,6 %",RED)]
for i, (g, th, sim, ec, col) in enumerate(comp):
    bg_r = LITE if i%2==0 else WITE
    for j, v in enumerate([g, th, sim, ec]):
        c = ct.cell(i+1, j)
        tcell(c, v, bg_c=bg_r, fc=col if j==3 else DT,
              sz=10.5, bold=(j==3))

# Bar chart
cd = ChartData()
cd.categories = ["Hd seuil","h coursier","h cuillère"]
cd.add_series("Écart (%)", (7.6, 36.6, 65.6))
s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED,
                   Inches(0.3), Inches(4.42), Inches(5.7), Inches(1.1), cd)

# Right analysis cards
ay = 0.88
for col, pct, title, text in [
    (GREEN,"7,6 %","Seuil Creager",
     "Excellent accord. La loi de Creager est robuste. L'écart reflète les effets de bords 3D."),
    (ORANG,"36,6 %","Coursier",
     "Valeur théorique : profondeur uniforme HEC-RAS. Simulation : section locale en régime non uniforme."),
    (RED,"65,6 %","Cuillère",
     "Valeur théorique : point de décollement strict. Simulation : hauteur dans la zone courbe."),
]:
    card(s, 6.15, ay, 3.55, 1.48)
    rect(s, 6.25, ay+0.1, 1.05, 0.48, fill=col, rnd=True)
    tx(s, pct, 6.25, ay+0.1, 1.05, 0.48, sz=14, fc=WITE, bold=True,
       al=PP_ALIGN.CENTER, fn="Cambria")
    tx(s, title, 7.38, ay+0.13, 2.22, 0.4, sz=12, fc=DT, bold=True, fn="Cambria")
    tx(s, text,  6.25, ay+0.65, 3.35, 0.74, sz=9, fc=MT)
    ay += 1.52
rect(s, 0.3, 5.3, 9.4, 0.24, fill=LITE, line=PRIM, rnd=True)
tx(s, "Les ordres de grandeur sont cohérents — les écarts reflètent les limites de la comparaison 1D/3D",
   0.42, 5.32, 9.18, 0.22, sz=8.5, fc=PRIM, italic=True, al=PP_ALIGN.CENTER)
note(s, "Slide 10 - Comparaison. 7,6% sur le seuil = excellent. Écarts coursier/cuillère = méthode différente, pas une erreur. (2 min)")

# ══ SLIDE 11 — Q1051 EXPLORATOIRE ════════════════════
s = dslide("Simulation exploratoire — Q = 1 051 m³/s",
           "Scénario post-surélévation (+32 % par rapport au débit millénal Q₁₀₀₀)")
stats = [
    (TEAL,"21 m/s","Vitesse max. sur le coursier","Sollicitation mécanique\nmaximale des bajoyers"),
    (PRIM,"Hd = 2,54 m","Charge sur le seuil Creager","Lue en ParaView\n(coordonnée Z surface libre)"),
    (RGB(0xB0,0x30,0x20),"Risque","Cavitation à la cuillère","Zone de sous-pression\ndétectée (P < 0)"),
]
for i, (col, val, label, subt) in enumerate(stats):
    sx = 0.3 + i*3.27
    rect(s, sx, 1.05, 3.05, 2.35, fill=col, rnd=True)
    tx(s, val,   sx+0.1, 1.15, 2.85, 0.75, sz=24, fc=WITE, bold=True,
       al=PP_ALIGN.CENTER, fn="Cambria")
    tx(s, label, sx+0.1, 1.93, 2.85, 0.45, sz=11, fc=LT, bold=True,
       al=PP_ALIGN.CENTER)
    tx(s, subt,  sx+0.1, 2.42, 2.85, 0.82, sz=9, fc=LT, italic=True,
       al=PP_ALIGN.CENTER)
tx(s, "Recommandations :", 0.4, 3.58, 5, 0.35, sz=13, fc=ACCN, bold=True, fn="Cambria")
recs = [
    "Renforcer les murs bajoyers du coursier (vitesses > 21 m/s)",
    "Protéger le lit de l'oued sur une portée minimale de 9 m en aval",
    "Traitement anti-cavitation à la base de la cuillère si surélévation retenue",
]
for i, r in enumerate(recs):
    tx(s, "►  "+r, 0.45, 4.02+i*0.44, 9.1, 0.4, sz=10.5, fc=WITE)
note(s, "Slide 11 - Q1051 exploratoire. Risque de cavitation = point critique. (1min30)")

# ══ SLIDE 12 — SYNTHÈSE ══════════════════════════════
s = lslide("Synthèse des résultats de simulation")
card(s, 0.3, 0.88, 9.4, 4.62)
tx(s, "Récapitulatif des grandeurs caractéristiques simulées",
   0.45, 0.95, 9.1, 0.38, sz=12, fc=PRIM, bold=True, fn="Cambria")
st = s.shapes.add_table(8, 4, Inches(0.38), Inches(1.42),
                         Inches(9.24), Inches(3.92)).table
for j, h in enumerate(["Grandeur","Q10 = 160,74 m³/s",
                        "Q100 = 507,54 m³/s","Q1000 = 794,93 m³/s"]):
    tcell(st.cell(0,j), h, bg_c=PRIM, fc=WITE, sz=10, bold=True)
syn_d = [
    ("Hd — seuil (m)",         "—",      "—",    "2,54"),
    ("h — coursier (m)",        "—",      "1,12", "1,12"),
    ("h — cuillère (m)",        "—",      "—",    "1,01"),
    ("R/y₀ — vérification",   "23,4 ✓","6,96 ✓","4,73 ✓"),
    ("Portée jet (m)",          "—",      "—",    "≈ 8,5"),
    ("Vitesse max. (m/s)",      "—",      "—",    "≈ 17,9"),
    ("Écart Hd vs théorie",     "—",      "—",    "7,6 %"),
]
for i, row in enumerate(syn_d):
    bg_r = LITE if i%2==0 else WITE
    for j, v in enumerate(row):
        c = st.cell(i+1, j)
        is_ok = "✓" in v
        tcell(c, v, bg_c=bg_r,
              fc=GREEN if is_ok else DT,
              sz=10, bold=is_ok)
note(s, "Slide 12 - Synthèse. R/y0 > 4 validé pour les 3 débits. (1 min)")

# ══ SLIDE 13 — CONCLUSION ════════════════════════════
s = lslide("Conclusion & Perspectives")
rect(s, 0.3, 0.88, 4.55, 4.62, fill=PRIM, rnd=True)
tx(s, "Ce que nous avons accompli", 0.45, 0.96, 4.25, 0.42,
   sz=12, fc=WITE, bold=True, fn="Cambria")
rect(s, 0.45, 1.42, 4.0, 0.03, fill=ACCN)
done = [
    "Dimensionnement complet (seuil, coursier, cuillère) Q10/Q100/Q1000",
    "Premier modèle CFD 3D OpenFOAM du Barrage Ahl Souss",
    "Validation croisée théorie/numérique — analyse des écarts",
    "Simulation exploratoire post-surélévation Q = 1 051 m³/s",
    "Critère R/y₀ > 4 vérifié pour les 3 débits de projet",
]
for i, item in enumerate(done):
    tx(s, "✓  "+item, 0.45, 1.52+i*0.6, 4.25, 0.54, sz=10.5, fc=LT)
card(s, 5.05, 0.88, 4.65, 2.12)
tx(s, "Limites de l'étude", 5.18, 0.96, 4.35, 0.38,
   sz=12, fc=RED, bold=True, fn="Cambria")
for i, l in enumerate(["Maillage fixe (pas d'adaptation automatique)",
                        "Modèle k-ε — moins précis en zones courbes",
                        "Absence de validation sur modèle physique"]):
    tx(s, "⚠  "+l, 5.18, 1.42+i*0.45, 4.42, 0.4, sz=10, fc=DT)
card(s, 5.05, 3.12, 4.65, 2.38)
tx(s, "Perspectives", 5.18, 3.2, 4.35, 0.38, sz=12, fc=TEAL, bold=True, fn="Cambria")
for i, p_ in enumerate(["Raffiner le maillage → réduire les écarts",
                         "Modèle RSM pour les écoulements courbes",
                         "Coupler avec transport sédimentaire",
                         "Valider par essai sur modèle physique"]):
    tx(s, "→  "+p_, 5.18, 3.65+i*0.47, 4.42, 0.42, sz=10, fc=DT)
note(s, "Slide 13 - Conclusion. Complémentarité des deux approches. (1 min)")

# ══ SLIDE 14 — QUESTIONS ══════════════════════════════
s = ns(); bg(s, DARK)
im(s, FIG+"figure4_location.png", 5.0, 0.4, 4.7, 5.0)
tx(s, "Merci de votre attention", 0.5, 1.05, 6.5, 0.75,
   sz=28, fc=WITE, bold=True, fn="Cambria")
rect(s, 0.5, 1.87, 5.5, 0.04, fill=ACCN)
tx(s, "Questions ?", 0.5, 2.05, 6.5, 1.0, sz=48, fc=ACCN, bold=True, fn="Cambria")
tx(s, "Hiba BASTOUT  &  Asmae BOULEMKHARJ",
   0.5, 3.3, 6, 0.4, sz=13, fc=LT)
tx(s, "École Hassania des Travaux Publics — Juin 2026",
   0.5, 3.75, 6, 0.35, sz=10, fc=LT, italic=True)
im(s, LOG+"ehtp-logo.png", 0.5, 4.65, 1.4, 0.72)
im(s, LOG+"cid-logo.jpg",  2.1, 4.65, 1.4, 0.72)
note(s, "Slide 14 - Questions. Durée totale ~18 min.")

# ── Save ──────────────────────────────────────────────
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
